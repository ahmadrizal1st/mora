import os
import sys
from PIL import Image
import docx
import openpyxl
from pptx import Presentation
import fitz  # PyMuPDF

# Ensure we can import from apps/ocr
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

def create_samples():
    sample_dir = os.path.join(os.path.dirname(__file__), "samples")
    os.makedirs(sample_dir, exist_ok=True)
    
    print(f"Generating samples in {sample_dir}...")

    # 1. TXT
    with open(os.path.join(sample_dir, "test.txt"), "w") as f:
        f.write("Hello World from TXT extraction test.")
    
    # 2. JPG/PNG
    img = Image.new('RGB', (200, 100), color=(255, 255, 255))
    from PIL import ImageDraw, ImageFont
    d = ImageDraw.Draw(img)
    d.text((10, 10), "Hello OCR Image", fill=(0, 0, 0))
    img.save(os.path.join(sample_dir, "test.jpg"))
    img.save(os.path.join(sample_dir, "test.png"))
    
    # 3. DOCX
    doc = docx.Document()
    doc.add_paragraph("Hello from DOCX extraction test.")
    doc.save(os.path.join(sample_dir, "test.docx"))
    
    # 4. XLSX
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Hello from XLSX extraction test."
    wb.save(os.path.join(sample_dir, "test.xlsx"))
    
    # 5. PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello from PPTX test"
    prs.save(os.path.join(sample_dir, "test.pptx"))
    
    # 6. PDF Text (Selectable)
    pdf_text = fitz.open()
    page = pdf_text.new_page()
    page.insert_text((50, 50), "This is a selectable PDF text for testing.")
    pdf_text.save(os.path.join(sample_dir, "test_text.pdf"))
    
    # 7. PDF Scan (Image-based PDF)
    pdf_scan = fitz.open()
    page = pdf_scan.new_page()
    page.insert_image(page.rect, filename=os.path.join(sample_dir, "test.jpg"))
    pdf_scan.save(os.path.join(sample_dir, "test_scan.pdf"))
    
    # 8. MP3 (Minimal valid-ish header)
    with open(os.path.join(sample_dir, "test.mp3"), "wb") as f:
        f.write(b"ID3\x03\x00\x00\x00\x00\x00\x23" + b"\x00" * 1024)

    # 9. Large File (101 MB) for 413 test
    with open(os.path.join(sample_dir, "large.txt"), "wb") as f:
        f.seek(101 * 1024 * 1024 - 1)
        f.write(b"\0")

    print("All samples generated successfully.")

if __name__ == "__main__":
    create_samples()
