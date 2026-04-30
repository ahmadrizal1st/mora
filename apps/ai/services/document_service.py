import docx
import openpyxl
from pptx import Presentation
from striprtf.striprtf import rtf_to_text

class DocumentService:
    @staticmethod
    def extract_text(file_path: str, mime_type: str) -> tuple[str, str]:
        """
        Extracts text from various document formats.
        Returns a tuple of (extracted_text, engine_used).
        """
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return DocumentService._extract_docx(file_path), "python-docx"
            
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return DocumentService._extract_xlsx(file_path), "openpyxl"
            
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return DocumentService._extract_pptx(file_path), "python-pptx"
            
        elif mime_type in ["text/plain", "text/csv"]:
            if file_path.lower().endswith('.csv') or mime_type == "text/csv":
                return DocumentService._extract_csv(file_path), "csv-builtin"
            return DocumentService._extract_txt(file_path), "builtin"
            
        elif mime_type in ["application/rtf", "text/rtf"]:
            return DocumentService._extract_rtf(file_path), "striprtf"
            
        return "", "none"

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        doc = docx.Document(file_path)
        text_content = []
        
        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Extract from tables (Crucial for bank statements/invoices)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                if row_text.replace("|", "").strip():
                    text_content.append(row_text)
                    
        return "\n".join(text_content)

    @staticmethod
    def _extract_xlsx(file_path: str) -> str:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text_content = []
        for sheet in wb.worksheets:
            text_content.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                # Use pipe separator to help LLM distinguish columns
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.replace("|", "").strip():
                    text_content.append(row_text)
        return "\n".join(text_content)

    @staticmethod
    def _extract_csv(file_path: str) -> str:
        import csv
        text_content = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                for row in reader:
                    row_text = " | ".join([cell.strip() for cell in row])
                    if row_text.replace("|", "").strip():
                        text_content.append(row_text)
        except Exception as e:
            return f"Error reading CSV: {str(e)}"
        return "\n".join(text_content)

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text_content = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
        return "\n".join(text_content)

    @staticmethod
    def _extract_txt(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    @staticmethod
    def _extract_rtf(file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
            return rtf_to_text(content)
